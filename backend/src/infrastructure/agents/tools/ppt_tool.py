"""
PowerPoint Tool — generates professional .pptx files from structured slide data.

Architecture:
  User → Agent → create_ppt tool
                    ↓
                SlideData (title + content + slide_type)
                    ↓
                python-pptx Presentation (with rich formatting + charts)
                    ↓
                storage/presentations/<name>_<uuid>.pptx
                    ↓
                Download URL

Design System:
  - Strict color palette (dark navy / slate / accent blue)
  - Consistent typography hierarchy (title → subtitle → body)
  - Uniform spacing & margins across all slides
  - Professional slide layouts: title, bullet, text, two-column, chart slides
  - Embedded charts (bar, line, pie) via Pillow

Supported Content Formatting:
  - **bold text** → bold runs
  - *italic text* → italic runs
  - - bullet points (with styled markers)
  - 1. numbered items
  - > key insight callouts (highlighted box)
  - | col1 | col2 | table syntax
  - ```chart ... ``` blocks (bar, line, pie)

Resource limits:
  - Max slides     : 30
  - Max file size   : 10 MB
"""

import os
import re
import io
import uuid
import math
import logging
import tempfile
from typing import Any, Dict, List, Optional, Tuple

from pydantic import BaseModel, Field
from sqlalchemy.orm import Session
from langchain_core.tools import StructuredTool

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

from src.infrastructure.database.models import User, Presentation, PresentationSlide

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Design System Constants
# ---------------------------------------------------------------------------

# Color palette
COLOR_BG_DARK     = RGBColor(0x1A, 0x1A, 0x2E)   # Deep navy background
COLOR_BG_SLIDE    = RGBColor(0xFF, 0xFF, 0xFF)     # White slide background
COLOR_TITLE       = RGBColor(0x1A, 0x1A, 0x2E)     # Dark navy for titles
COLOR_SUBTITLE    = RGBColor(0x4A, 0x4A, 0x6A)     # Slate for subtitles
COLOR_BODY        = RGBColor(0x33, 0x33, 0x44)     # Dark gray for body text
COLOR_ACCENT      = RGBColor(0x2D, 0x7D, 0xD2)     # Accent blue for highlights
COLOR_BULLET      = RGBColor(0x2D, 0x7D, 0xD2)     # Accent blue for bullet markers
COLOR_DIVIDER     = RGBColor(0xE0, 0xE4, 0xE8)     # Light gray divider
COLOR_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG    = RGBColor(0xF0, 0xF4, 0xF8)     # Light blue-gray for callout bg
COLOR_TABLE_HDR   = RGBColor(0x1A, 0x1A, 0x2E)     # Navy header
COLOR_TABLE_ALT   = RGBColor(0xF0, 0xF4, 0xF8)     # Alternating row
COLOR_MUTED       = RGBColor(0x6B, 0x70, 0x7B)     # Muted gray for footnotes
COLOR_INSIGHT_BG  = RGBColor(0xEB, 0xF3, 0xFD)     # Light blue for insight callouts

# Chart color palette (RGB tuples for Pillow)
CHART_COLORS = [
    (0x2D, 0x7D, 0xD2),   # Blue
    (0xE8, 0x6C, 0x50),   # Coral-red
    (0x2E, 0xCC, 0x71),   # Green
    (0xF3, 0x9C, 0x12),   # Amber
    (0x9B, 0x59, 0xB6),   # Purple
    (0x1A, 0xBC, 0x9C),   # Teal
    (0xE7, 0x4C, 0x3C),   # Red
    (0x34, 0x98, 0xDB),   # Light blue
    (0xF1, 0xC4, 0x0F),   # Yellow
    (0x2C, 0x3E, 0x50),   # Dark slate
]

# Typography
FONT_FAMILY       = "Calibri"
FONT_TITLE_SIZE   = Pt(32)
FONT_HEADING_SIZE = Pt(24)
FONT_BODY_SIZE    = Pt(16)
FONT_SMALL_SIZE   = Pt(12)

# Layout (16:9 widescreen)
SLIDE_WIDTH       = Inches(13.333)
SLIDE_HEIGHT      = Inches(7.5)
MARGIN_LEFT       = Inches(0.8)
MARGIN_TOP        = Inches(0.6)
CONTENT_WIDTH     = Inches(11.7)
CONTENT_TOP       = Inches(1.6)
CONTENT_HEIGHT    = Inches(5.2)

# Limits
MAX_SLIDES = 30
MAX_FILE_SIZE_BYTES = 10 * 1024 * 1024

# ---------------------------------------------------------------------------
# Input Schemas
# ---------------------------------------------------------------------------


class PptSlideInput(BaseModel):
    """Represents one slide in the presentation."""
    title: str = Field(description="Title for this slide.")
    content: str = Field(
        description=(
            "Content for this slide. Supports: '- ' for bullet points, "
            "'1. ' for numbered items, '> ' for insight callouts, "
            "'**bold**', '*italic*', '| col | col |' for tables, "
            "and ```chart blocks for embedded charts."
        )
    )
    slide_type: str = Field(
        default="bullet",
        description=(
            "Type of slide: 'bullet' for bullet-point slides, "
            "'text' for paragraph slides, 'title' for section divider, "
            "'two_column' for two-column layout (separate columns with |||)."
        )
    )


class CreatePptInput(BaseModel):
    """Input schema for the create_ppt tool."""
    title: str = Field(description="Title of the PowerPoint presentation.")
    slides: List[PptSlideInput] = Field(
        description=(
            "List of slides. Each has 'title', 'content' (supports markdown), "
            "and optional 'slide_type' ('bullet', 'text', 'title', 'two_column')."
        )
    )
    base_url: str = Field(
        description="Base URL of the backend server. Injected automatically."
    )


# ---------------------------------------------------------------------------
# Chart rendering (using Pillow)
# ---------------------------------------------------------------------------


def _get_text_size(draw, text: str, font) -> Tuple[int, int]:
    """Get text bounding box size."""
    try:
        bbox = draw.textbbox((0, 0), text, font=font)
        return bbox[2] - bbox[0], bbox[3] - bbox[1]
    except AttributeError:
        return draw.textsize(text, font=font)


def _get_font(size: int, bold: bool = False):
    """Try to load a good font."""
    font_paths = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ]
    for path in font_paths:
        try:
            return ImageFont.truetype(path, size)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()


def _render_bar_chart(title: str, labels: List[str], values: List[float],
                      width: int = 900, height: int = 500) -> Image.Image:
    """Render a bar chart."""
    img = Image.new('RGB', (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    title_font = _get_font(20, bold=True)
    label_font = _get_font(13)
    value_font = _get_font(12, bold=True)

    tw, th = _get_text_size(draw, title, title_font)
    draw.text(((width - tw) // 2, 12), title, fill=(0x1A, 0x1A, 0x2E), font=title_font)

    ml, mr, mt, mb = 80, 40, 55, 75
    cw = width - ml - mr
    ch = height - mt - mb
    if not values:
        return img

    max_val = max(values) if max(values) > 0 else 1
    n = len(labels)
    bar_w = max(25, min(65, cw // (n * 2)))
    spacing = (cw - bar_w * n) / (n + 1)

    for i in range(6):
        y = mt + ch - (i * ch / 5)
        draw.line([(ml, int(y)), (width - mr, int(y))], fill=(0xE8, 0xE8, 0xE8), width=1)
        val_str = f"{(i * max_val / 5):,.0f}"
        vw, vh = _get_text_size(draw, val_str, label_font)
        draw.text((ml - vw - 6, int(y) - vh // 2), val_str, fill=(0x90, 0x90, 0x90), font=label_font)

    for i, (label, val) in enumerate(zip(labels, values)):
        x = ml + spacing + i * (bar_w + spacing)
        bh = (val / max_val) * ch
        yt = mt + ch - bh
        color = CHART_COLORS[i % len(CHART_COLORS)]
        draw.rectangle([int(x), int(yt), int(x + bar_w), int(mt + ch)], fill=color)
        vs = f"{val:,.0f}" if val == int(val) else f"{val:,.1f}"
        vw, vh = _get_text_size(draw, vs, value_font)
        draw.text((int(x + bar_w / 2 - vw / 2), int(yt - vh - 3)), vs, fill=color, font=value_font)
        lw, lh = _get_text_size(draw, label, label_font)
        draw.text((int(x + bar_w / 2 - lw / 2), int(mt + ch + 6)), label, fill=(0x33, 0x33, 0x44), font=label_font)

    draw.line([(ml, mt), (ml, mt + ch)], fill=(0x33, 0x33, 0x44), width=2)
    draw.line([(ml, mt + ch), (width - mr, mt + ch)], fill=(0x33, 0x33, 0x44), width=2)
    return img


def _render_line_chart(title: str, labels: List[str], values: List[float],
                       width: int = 900, height: int = 500) -> Image.Image:
    """Render a line chart."""
    img = Image.new('RGB', (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    title_font = _get_font(20, bold=True)
    label_font = _get_font(13)
    value_font = _get_font(11, bold=True)

    tw, th = _get_text_size(draw, title, title_font)
    draw.text(((width - tw) // 2, 12), title, fill=(0x1A, 0x1A, 0x2E), font=title_font)

    ml, mr, mt, mb = 80, 40, 55, 75
    cw = width - ml - mr
    ch = height - mt - mb
    if not values or len(values) < 2:
        return img

    max_v = max(values) if max(values) > 0 else 1
    min_v = min(values)
    vr = max_v - min_v if max_v != min_v else 1
    n = len(labels)

    for i in range(6):
        y = mt + ch - (i * ch / 5)
        draw.line([(ml, int(y)), (width - mr, int(y))], fill=(0xE8, 0xE8, 0xE8), width=1)

    points = []
    step = cw / (n - 1) if n > 1 else cw
    for i, val in enumerate(values):
        x = ml + i * step
        y = mt + ch - ((val - min_v) / vr) * ch
        points.append((int(x), int(y)))

    color = CHART_COLORS[0]
    overlay = Image.new('RGBA', (width, height), (255, 255, 255, 0))
    od = ImageDraw.Draw(overlay)
    area = points + [(points[-1][0], mt + ch), (points[0][0], mt + ch)]
    od.polygon(area, fill=(color[0], color[1], color[2], 35))
    img.paste(Image.alpha_composite(Image.new('RGBA', img.size, (255, 255, 255, 255)), overlay).convert('RGB'))
    draw = ImageDraw.Draw(img)

    for i in range(6):
        y = mt + ch - (i * ch / 5)
        draw.line([(ml, int(y)), (width - mr, int(y))], fill=(0xE8, 0xE8, 0xE8), width=1)

    for i in range(len(points) - 1):
        draw.line([points[i], points[i + 1]], fill=color, width=3)
    for i, (px, py) in enumerate(points):
        draw.ellipse([px - 5, py - 5, px + 5, py + 5], fill=color, outline=(255, 255, 255))
        vs = f"{values[i]:,.0f}" if values[i] == int(values[i]) else f"{values[i]:,.1f}"
        vw, vh = _get_text_size(draw, vs, value_font)
        draw.text((px - vw // 2, py - vh - 8), vs, fill=color, font=value_font)
    for i, label in enumerate(labels):
        x = ml + i * step
        lw, lh = _get_text_size(draw, label, label_font)
        draw.text((int(x - lw / 2), mt + ch + 6), label, fill=(0x33, 0x33, 0x44), font=label_font)

    draw.line([(ml, mt), (ml, mt + ch)], fill=(0x33, 0x33, 0x44), width=2)
    draw.line([(ml, mt + ch), (width - mr, mt + ch)], fill=(0x33, 0x33, 0x44), width=2)
    return img


def _render_pie_chart(title: str, labels: List[str], values: List[float],
                      width: int = 800, height: int = 500) -> Image.Image:
    """Render a pie chart."""
    img = Image.new('RGB', (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    title_font = _get_font(20, bold=True)
    label_font = _get_font(13)

    tw, th = _get_text_size(draw, title, title_font)
    draw.text(((width - tw) // 2, 12), title, fill=(0x1A, 0x1A, 0x2E), font=title_font)

    if not values or sum(values) == 0:
        return img
    total = sum(values)
    cx, cy = width // 2 - 70, height // 2 + 15
    radius = min(width, height) // 2 - 80

    angle = -90
    for i, (label, val) in enumerate(zip(labels, values)):
        sweep = (val / total) * 360
        color = CHART_COLORS[i % len(CHART_COLORS)]
        draw.pieslice([cx - radius, cy - radius, cx + radius, cy + radius],
                      angle, angle + sweep, fill=color, outline=(255, 255, 255), width=2)
        angle += sweep

    lx = cx + radius + 35
    ly = cy - (len(labels) * 26) // 2
    for i, (label, val) in enumerate(zip(labels, values)):
        color = CHART_COLORS[i % len(CHART_COLORS)]
        y = ly + i * 26
        draw.rectangle([lx, y, lx + 14, y + 14], fill=color)
        draw.text((lx + 22, y), f"{label} ({(val / total) * 100:.1f}%)",
                  fill=(0x33, 0x33, 0x44), font=label_font)
    return img


def _parse_chart_block(chart_text: str) -> Optional[Dict[str, Any]]:
    """Parse a chart definition block."""
    result = {"type": "bar", "title": "Chart", "labels": [], "values": []}
    in_data = False
    for line in chart_text.strip().split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.lower().startswith("type:"):
            result["type"] = stripped.split(":", 1)[1].strip().lower()
            in_data = False
        elif stripped.lower().startswith("title:"):
            result["title"] = stripped.split(":", 1)[1].strip()
            in_data = False
        elif stripped.lower() == "data:":
            in_data = True
        elif in_data and ":" in stripped:
            parts = stripped.rsplit(":", 1)
            if len(parts) == 2:
                try:
                    value = float(parts[1].strip().replace(",", "").replace("$", "").replace("%", ""))
                    result["labels"].append(parts[0].strip())
                    result["values"].append(value)
                except ValueError:
                    continue
    return result if result["labels"] else None


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set the background color of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_slide_number(slide, slide_num: int, total: int) -> None:
    """Add slide number to bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total}"
    run.font.size = Pt(9)
    run.font.color.rgb = COLOR_MUTED
    run.font.name = FONT_FAMILY


def _parse_inline_runs(paragraph, text: str, font_size=FONT_BODY_SIZE,
                       font_color=COLOR_BODY) -> None:
    """Parse inline markdown (**bold**, *italic*) and add formatted runs."""
    pattern = re.compile(
        r'(\*\*\*(.+?)\*\*\*)'
        r'|(\*\*(.+?)\*\*)'
        r'|(\*(.+?)\*)'
    )
    last_end = 0
    for match in pattern.finditer(text):
        if match.start() > last_end:
            run = paragraph.add_run()
            run.text = text[last_end:match.start()]
            run.font.size = font_size
            run.font.color.rgb = font_color
            run.font.name = FONT_FAMILY

        if match.group(2):  # ***bold italic***
            run = paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
            run.font.italic = True
        elif match.group(4):  # **bold**
            run = paragraph.add_run()
            run.text = match.group(4)
            run.font.bold = True
        elif match.group(6):  # *italic*
            run = paragraph.add_run()
            run.text = match.group(6)
            run.font.italic = True

        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.name = FONT_FAMILY
        last_end = match.end()

    if last_end < len(text):
        run = paragraph.add_run()
        run.text = text[last_end:]
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.name = FONT_FAMILY


def _add_table_to_slide(slide, lines: List[str]) -> None:
    """Parse markdown table and add to slide."""
    rows_data = []
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|"):
            continue
        if re.match(r'^\|[\s\-:]+(\|[\s\-:]+)*\|?$', stripped):
            continue
        cells = [c.strip() for c in stripped.split("|")]
        if cells and cells[0] == "":
            cells = cells[1:]
        if cells and cells[-1] == "":
            cells = cells[:-1]
        if cells:
            rows_data.append(cells)

    if not rows_data:
        return

    num_cols = max(len(row) for row in rows_data)
    for row in rows_data:
        while len(row) < num_cols:
            row.append("")

    # Table dimensions
    table_width = CONTENT_WIDTH
    col_width = int(table_width / num_cols)
    row_height = Inches(0.4)

    table_shape = slide.shapes.add_table(
        len(rows_data), num_cols,
        MARGIN_LEFT, Inches(3.2),  # Position below title area
        table_width, row_height * len(rows_data)
    )
    table = table_shape.table

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            run = p.add_run()
            run.text = cell_text
            run.font.name = FONT_FAMILY

            if row_idx == 0:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE
                # Header background
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = COLOR_TABLE_HDR
            else:
                run.font.size = Pt(10)
                run.font.color.rgb = COLOR_BODY
                if row_idx % 2 == 0:
                    cell_fill = cell.fill
                    cell_fill.solid()
                    cell_fill.fore_color.rgb = COLOR_TABLE_ALT


def _add_chart_to_slide(slide, chart_data: Dict[str, Any]) -> None:
    """Render a chart and add as image to slide."""
    if not PIL_AVAILABLE:
        return

    chart_type = chart_data.get("type", "bar")
    title = chart_data.get("title", "Chart")
    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])

    try:
        if chart_type == "line":
            chart_img = _render_line_chart(title, labels, values)
        elif chart_type == "pie":
            chart_img = _render_pie_chart(title, labels, values)
        else:
            chart_img = _render_bar_chart(title, labels, values)

        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            chart_img.save(tmp, format="PNG", quality=95)
            tmp_path = tmp.name

        # Center chart on slide
        img_width = Inches(8.5)
        img_left = (SLIDE_WIDTH - img_width) // 2
        slide.shapes.add_picture(tmp_path, img_left, Inches(2.0), width=img_width)

        try:
            os.remove(tmp_path)
        except OSError:
            pass
    except Exception as e:
        logger.warning(f"Failed to render chart in PPT: {e}")


def _add_insight_callout(slide, text: str, top: Inches) -> None:
    """Add a styled insight/callout box to a slide."""
    # Background shape
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        MARGIN_LEFT, top,
        CONTENT_WIDTH, Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_INSIGHT_BG
    shape.line.fill.background()

    # Left accent bar
    accent = slide.shapes.add_shape(
        1, MARGIN_LEFT, top,
        Inches(0.06), Inches(0.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

    # Text
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.3), top + Inches(0.1),
        CONTENT_WIDTH - Inches(0.5), Inches(0.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x2D, 0x5F, 0x8A)
    run.font.name = FONT_FAMILY


def _build_title_slide(prs, title: str, subtitle: str = "") -> None:
    """Build a cover/title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_BG_DARK)

    # Title
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.3), CONTENT_WIDTH, Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.color.rgb = COLOR_WHITE
    run.font.bold = True
    run.font.name = FONT_FAMILY

    # Accent divider
    shape = slide.shapes.add_shape(
        1, Inches(5.0), Inches(3.9),
        Inches(3.3), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(4.2), CONTENT_WIDTH, Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        run2.font.name = FONT_FAMILY


def _build_section_divider(prs, title: str, subtitle: str = "") -> None:
    """Build a section divider slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_BG_DARK)

    # Section number / title
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.8), CONTENT_WIDTH, Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.color.rgb = COLOR_WHITE
    run.font.bold = True
    run.font.name = FONT_FAMILY

    # Accent line under title
    shape = slide.shapes.add_shape(
        1, MARGIN_LEFT, Inches(3.95),
        Inches(3.0), Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(4.2), CONTENT_WIDTH, Inches(0.8)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        run2.font.name = FONT_FAMILY


def _build_content_slide(prs, title: str, content: str, slide_type: str = "bullet",
                         slide_num: int = 0, total_slides: int = 0) -> None:
    """Build a content slide with rich formatting support."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, COLOR_BG_SLIDE)

    # Slide title
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, CONTENT_WIDTH, Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = FONT_HEADING_SIZE
    run.font.color.rgb = COLOR_TITLE
    run.font.bold = True
    run.font.name = FONT_FAMILY

    # Accent divider under title
    shape = slide.shapes.add_shape(
        1, MARGIN_LEFT, Inches(1.3),
        Inches(2.0), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()

    # Handle two-column layout
    if slide_type == "two_column" and "|||" in content:
        cols = content.split("|||", 1)
        _render_content_in_textbox(slide, cols[0].strip(),
                                   MARGIN_LEFT, CONTENT_TOP,
                                   Inches(5.5), CONTENT_HEIGHT)
        _render_content_in_textbox(slide, cols[1].strip(),
                                   Inches(7.0), CONTENT_TOP,
                                   Inches(5.5), CONTENT_HEIGHT)
    else:
        # Parse for charts and tables that need special handling
        lines = content.split("\n")
        has_chart = False
        has_table = False
        chart_data = None
        table_lines_collected = []

        # Pre-scan for charts and tables
        i = 0
        regular_lines = []
        while i < len(lines):
            line = lines[i].strip()
            if line.startswith("```chart"):
                i += 1
                chart_block = []
                while i < len(lines) and not lines[i].strip().startswith("```"):
                    chart_block.append(lines[i])
                    i += 1
                if i < len(lines):
                    i += 1
                chart_data = _parse_chart_block("\n".join(chart_block))
                has_chart = True
            elif line.startswith("|"):
                while i < len(lines) and lines[i].strip().startswith("|"):
                    table_lines_collected.append(lines[i])
                    i += 1
                has_table = True
            else:
                regular_lines.append(lines[i])
                i += 1

        # Render regular content in textbox
        if regular_lines:
            content_text = "\n".join(regular_lines)
            content_top = CONTENT_TOP
            content_h = CONTENT_HEIGHT

            # Adjust if chart or table will follow
            if has_chart or has_table:
                content_h = Inches(2.5)

            _render_content_in_textbox(slide, content_text,
                                       MARGIN_LEFT, content_top,
                                       CONTENT_WIDTH, content_h)

        # Add chart if present
        if has_chart and chart_data:
            _add_chart_to_slide(slide, chart_data)

        # Add table if present
        if has_table and table_lines_collected:
            _add_table_to_slide(slide, table_lines_collected)

    # Slide number
    if slide_num > 0 and total_slides > 0:
        _add_slide_number(slide, slide_num, total_slides)


def _render_content_in_textbox(slide, content: str, left, top, width, height) -> None:
    """Render content with inline formatting into a textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = content.split("\n")
    first = True

    for line in lines:
        line_stripped = line.strip()
        if not line_stripped:
            continue

        # Skip chart blocks and table lines (handled separately)
        if line_stripped.startswith("```") or line_stripped.startswith("|"):
            continue

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        # > Insight callout (render as highlighted text)
        if line_stripped.startswith("> "):
            text = line_stripped[2:].strip()
            run = p.add_run()
            run.text = "💡 "
            run.font.size = Pt(14)
            run.font.name = FONT_FAMILY
            run2 = p.add_run()
            run2.text = text
            run2.font.size = Pt(14)
            run2.font.italic = True
            run2.font.color.rgb = RGBColor(0x2D, 0x5F, 0x8A)
            run2.font.name = FONT_FAMILY
            p.space_after = Pt(10)
            p.space_before = Pt(6)
            continue

        # - Bullet point
        if line_stripped.startswith("- ") or line_stripped.startswith("* "):
            text = line_stripped[2:].strip()
            # Bullet marker
            run = p.add_run()
            run.text = "●  "
            run.font.size = Pt(10)
            run.font.color.rgb = COLOR_BULLET
            run.font.name = FONT_FAMILY
            # Content with inline formatting
            _parse_inline_runs(p, text, FONT_BODY_SIZE, COLOR_BODY)
            p.space_after = Pt(8)
            p.space_before = Pt(2)
            continue

        # 1. Numbered item
        if re.match(r'^\d+\.\s+', line_stripped):
            num_match = re.match(r'^(\d+)\.\s+', line_stripped)
            num = num_match.group(1) if num_match else "1"
            text = re.sub(r'^\d+\.\s+', '', line_stripped)
            run = p.add_run()
            run.text = f"{num}.  "
            run.font.size = FONT_BODY_SIZE
            run.font.color.rgb = COLOR_ACCENT
            run.font.bold = True
            run.font.name = FONT_FAMILY
            _parse_inline_runs(p, text, FONT_BODY_SIZE, COLOR_BODY)
            p.space_after = Pt(8)
            p.space_before = Pt(2)
            continue

        # Regular text with inline formatting
        _parse_inline_runs(p, line_stripped, FONT_BODY_SIZE, COLOR_BODY)
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        p.alignment = PP_ALIGN.LEFT


# ---------------------------------------------------------------------------
# Core logic
# ---------------------------------------------------------------------------


async def create_ppt_db(
    input_data: CreatePptInput,
    sql_db: Session,
    user_id: str,
    conversation_id: str,
) -> str:
    """Generate a .pptx file from structured slide data and persist a DB record."""
    try:
        # 1. Validate user / org
        user = sql_db.query(User).filter(User.id == user_id).first()
        if not user or not user.org_id:
            return "❌ User or Organization not found."

        # 2. Validate slide count
        if len(input_data.slides) > MAX_SLIDES:
            return f"❌ Too many slides ({len(input_data.slides)}). Maximum is {MAX_SLIDES}."

        # 3. Prepare output path
        sanitized_title = re.sub(r'[^\w\s-]', '', input_data.title).strip().replace(' ', '_')
        if not sanitized_title:
            sanitized_title = "presentation"

        unique_suffix = str(uuid.uuid4())[:8]
        filename = f"{sanitized_title}_{unique_suffix}.pptx"

        storage_rel = f"storage/presentations/{filename}"
        os.makedirs("storage/presentations", exist_ok=True)

        # 4. Build PowerPoint
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # Count total slides for numbering
        total_content_slides = sum(1 for s in input_data.slides if s.slide_type != "title")

        # Cover slide
        _build_title_slide(prs, input_data.title)

        # Content slides
        content_num = 0
        for slide_data in input_data.slides:
            if slide_data.slide_type == "title":
                _build_section_divider(prs, slide_data.title, slide_data.content)
            else:
                content_num += 1
                _build_content_slide(
                    prs, slide_data.title, slide_data.content,
                    slide_data.slide_type,
                    slide_num=content_num, total_slides=total_content_slides
                )

        # 5. Save
        prs.save(storage_rel)

        # 6. Check file size
        file_size = os.path.getsize(storage_rel)
        if file_size > MAX_FILE_SIZE_BYTES:
            os.remove(storage_rel)
            return (
                f"❌ Generated file is {file_size / (1024 * 1024):.1f} MB which "
                f"exceeds the 10 MB limit. Please reduce the number of slides."
            )

        # 7. Build public URL
        base = input_data.base_url.rstrip("/")
        file_url = f"{base}/{storage_rel}"

        # 8. Persist DB record
        record = Presentation(
            conversation_id=conversation_id,
            org_id=user.org_id,
            title=input_data.title,
            file_path=storage_rel,
            file_url=file_url,
        )
        sql_db.add(record)
        sql_db.flush()

        for i, slide_in in enumerate(input_data.slides):
            slide_record = PresentationSlide(
                presentation_id=record.id,
                title=slide_in.title,
                content=slide_in.content,
                slide_type=slide_in.slide_type,
                order=i + 1,
            )
            sql_db.add(slide_record)

        sql_db.commit()
        sql_db.refresh(record)

        total_slides = len(input_data.slides) + 1  # +1 for cover
        return (
            f"✅ Presentation **'{record.title}'** created successfully!\n"
            f"📊 Total slides: {total_slides} (including cover)\n"
            f"📥 **Download link:** {file_url}"
        )

    except Exception as e:
        sql_db.rollback()
        logger.error("Error creating presentation: %s", e, exc_info=True)
        return f"❌ Error creating presentation: {str(e)}"


def get_ppt_link_db(sql_db: Session, conversation_id: str) -> str:
    """Return the download link for the latest presentation in this conversation."""
    try:
        record = (
            sql_db.query(Presentation)
            .filter(Presentation.conversation_id == conversation_id)
            .order_by(Presentation.created_at.desc())
            .first()
        )
        if not record:
            return "📋 No presentation found for this conversation. Use 'create_ppt' to generate one."
        if not record.file_url:
            return "📋 Presentation record found but the file has not been generated yet."
        return (
            f"📥 **Presentation:** '{record.title}'\n"
            f"**Download:** {record.file_url}"
        )
    except Exception as e:
        return f"❌ Error retrieving presentation link: {str(e)}"


# ---------------------------------------------------------------------------
# StructuredTool integration (LangChain)
# ---------------------------------------------------------------------------


def ppt_generation_tool(
    sql_db: Session,
    user_id: str,
    conversation_id: Optional[str] = None,
    base_url: str = "http://localhost:8000",
) -> List[StructuredTool]:
    """Returns LangChain StructuredTools for PowerPoint generation."""
    if not conversation_id:
        logger.warning("ppt_generation_tool called without conversation_id.")
        return []

    async def create_ppt(
        title: str,
        slides: List[Dict[str, Any]],
    ) -> str:
        """
        Generate a professional PowerPoint (.pptx) presentation from structured
        slide data and return a download link.

        The presentation will include a cover slide, slide numbers, professional
        design with accent colors, and support for rich content.

        CONTENT FORMATTING — the 'content' field supports:
          - **bold text** → bold
          - *italic text* → italic
          - - item → styled bullet point (blue marker)
          - 1. item → styled numbered list
          - > text → insight callout (highlighted, italic)
          - | Col1 | Col2 | → formatted table with styled header
          - ```chart ... ``` → embedded chart (bar, line, or pie)
          - ||| → column separator (use with slide_type="two_column")

        CHART FORMAT (inside content):
          ```chart
          type: bar
          title: Revenue by Quarter
          data:
            Q1: 180
            Q2: 210
            Q3: 245
            Q4: 285
          ```
          Supported types: bar, line, pie

        SLIDE TYPES:
          - "bullet" (default): Standard content slide with bullet points
          - "text": Paragraph-style content slide
          - "title": Section divider slide (dark background)
          - "two_column": Two-column layout (separate columns with |||)

        LIMITS:
          - Max 30 slides
          - Max 10 MB file size
        """
        parsed_slides = [
            PptSlideInput(**s) if isinstance(s, dict) else s for s in slides
        ]
        input_obj = CreatePptInput(
            title=title,
            slides=parsed_slides,
            base_url=base_url,
        )
        return await create_ppt_db(input_obj, sql_db, user_id, conversation_id)

    def get_ppt_link() -> str:
        """Return the download link for the most recently created presentation."""
        return get_ppt_link_db(sql_db, conversation_id)

    class _CreateInput(BaseModel):
        title: str = Field(description="Title of the PowerPoint presentation.")
        slides: List[Dict[str, Any]] = Field(
            description=(
                "List of slide objects. Each must have 'title' (str) and "
                "'content' (str, supports: **bold**, *italic*, - bullets, "
                "> insights, | tables |, ```chart blocks). "
                "Optional 'slide_type': 'bullet', 'text', 'title', 'two_column'. "
                "Max 30 slides."
            )
        )

    return [
        StructuredTool.from_function(
            coroutine=create_ppt,
            name="create_ppt",
            description=(
                "Generate a professional PowerPoint (.pptx) presentation with cover slide, "
                "slide numbers, charts, tables, and rich formatting. Pass 'title' and 'slides' "
                "(list of {title, content} objects). Content supports: **bold**, *italic*, "
                "- bullets, > insights, | tables |, ```chart blocks. "
                "slide_type options: 'bullet', 'text', 'title' (divider), 'two_column'. "
                "Returns a download link. Max 30 slides, 10 MB."
            ),
            args_schema=_CreateInput,
        ),
        StructuredTool.from_function(
            func=get_ppt_link,
            name="get_ppt_link",
            description=(
                "Retrieve the download link for the most recent PowerPoint "
                "created in this conversation."
            ),
            args_schema=None,
        ),
    ]
